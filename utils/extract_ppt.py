# from pptx import Presentation
# from PIL import Image
# import base64
# import io

# def extract_from_ppt(file_path):
#     pres = Presentation(file_path)
#     full_text = ""
#     images_b64 = []

#     for slide in pres.slides:
#         for shape in slide.shapes:
#             if shape.has_text_frame:
#                 full_text += shape.text + "\n"
#             if shape.shape_type == 13:
#                 image = shape.image.blob
#                 img = Image.open(io.BytesIO(image))
#                 buff = io.BytesIO()
#                 img.save(buff, format="PNG")
#                 images_b64.append(base64.b64encode(buff.getvalue()).decode())

#     return full_text, images_b64
from pptx import Presentation
from PIL import Image
import base64
import io

def extract_from_ppt(file_path):
    prs = Presentation(file_path)
    slides_data = []
    
    for slide_idx, slide in enumerate(prs.slides, 1):
        # --- Extract text from this slide ---
        slide_text = ""
        for shape in slide.shapes:
            if shape.has_text_frame:
                slide_text += shape.text + "\n"
        
        # --- Extract and convert slide to image ---
        slide_image_b64 = None
        
        # Try to extract embedded images from shapes
        for shape in slide.shapes:
            if shape.shape_type == 13:  # picture (MSO_SHAPE_TYPE.PICTURE)
                try:
                    image_blob = shape.image.blob
                    img = Image.open(io.BytesIO(image_blob))
                    
                    # Convert to RGB if needed (handles RGBA, P, etc.)
                    if img.mode not in ('RGB', 'L'):
                        img = img.convert('RGB')
                    
                    # Save as PNG with proper format
                    buff = io.BytesIO()
                    img.save(buff, format="PNG")
                    buff.seek(0)
                    slide_image_b64 = base64.b64encode(buff.getvalue()).decode('utf-8')
                    break  # Use first valid image found
                    
                except Exception as e:
                    print(f"Warning: Could not extract image from slide {slide_idx}: {e}")
                    continue
        
        # Store slide data (text, image)
        slides_data.append((slide_text.strip(), slide_image_b64))
    
    return None, slides_data